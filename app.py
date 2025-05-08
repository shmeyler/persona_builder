import streamlit as st
import pandas as pd
import io
import docx
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, TimeoutError
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError
from google.cloud import vision
import openai

# === Setup & Auth ===
st.title("📂 Persona Builder — Multi-Format + GPT")

creds = service_account.Credentials.from_service_account_info(st.secrets["gcp"])
drive_service = build("drive", "v3", credentials=creds)
vision_client = vision.ImageAnnotatorClient(credentials=creds)
openai.api_key = st.secrets["openai_api_key"]

FOLDER_ID = "1QBUwWvuaLvJrie3cblt8d4ch9cyaogWg"

# === List Files Recursively ===
def list_all_files(folder_id):
    query = f"'{folder_id}' in parents"
    results = drive_service.files().list(q=query, pageSize=100, fields="files(id, name, mimeType)").execute()
    items = results.get("files", [])
    all_files = []
    for item in items:
        if item["mimeType"] == "application/vnd.google-apps.folder":
            st.write(f"📁 Entering folder: {item['name']}")
            all_files.extend(list_all_files(item["id"]))
        else:
            all_files.append(item)
    return all_files

# === Timeout-Safe Parsers ===
def safe_parse_pptx(fh):
    def parse():
        prs = Presentation(fh)
        return "\n".join(
            shape.text.strip()
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    with ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(parse)
        return future.result(timeout=5)

def safe_ocr_png(image_bytes):
    def ocr():
        image = vision.Image(content=image_bytes)
        response = vision_client.text_detection(image=image)
        annotations = response.text_annotations
        return annotations[0].description.strip() if annotations else ""
    with ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(ocr)
        return future.result(timeout=5)

# === Main Parser Logic ===
st.write("🔍 Scanning Google Drive folder...")
files = list_all_files(FOLDER_ID)
all_texts = []

if not files:
    st.warning("No files found.")
else:
    st.success(f"✅ Found {len(files)} files.")

    for file in files:
        file_id = file["id"]
        file_name = file["name"]
        mime = file["mimeType"]

        st.write(f"📄 Processing: {file_name} ({mime})")

        try:
            fh = io.BytesIO()

            if mime.startswith("application/vnd.google-apps"):
                export_mime = {
                    "application/vnd.google-apps.document": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/vnd.google-apps.spreadsheet": "text/csv",
                    "application/vnd.google-apps.presentation": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                }.get(mime)
                if not export_mime:
                    st.warning(f"⏭️ Skipping unsupported Google-native file: {file_name}")
                    continue
                request = drive_service.files().export_media(fileId=file_id, mimeType=export_mime)
            else:
                request = drive_service.files().get_media(fileId=file_id)

            downloader = MediaIoBaseDownload(fh, request)
            while not downloader.next_chunk()[1]:
                pass
            fh.seek(0)

            text = ""
            if file_name.endswith(".csv"):
                df = pd.read_csv(fh)
                text = df.to_csv(index=False)

            elif file_name.endswith(".xlsx"):
                df = pd.read_excel(fh)
                text = df.to_csv(index=False)

            elif file_name.endswith(".docx"):
                doc = docx.Document(fh)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

            elif file_name.endswith(".pptx") or "presentation" in mime:
                try:
                    text = safe_parse_pptx(fh)
                except TimeoutError:
                    st.warning(f"⏱️ Skipped {file_name} (pptx timeout)")
                    continue

            elif file_name.lower().endswith(".png") and mime.startswith("image/"):
                try:
                    st.write(f"🖼️ Running OCR on {file_name}...")
                    image_content = fh.read()
                    text = safe_ocr_png(image_content)
                    if text:
                        st.success(f"✅ OCR complete: {file_name}")
                    else:
                        st.warning(f"⚪ No text found in {file_name}")
                        continue
                except TimeoutError:
                    st.warning(f"⏱️ Skipped {file_name} (OCR timeout)")
                    continue

            else:
                st.warning(f"⏭️ Skipping unsupported type: {file_name}")
                continue

            if text.strip():
                all_texts.append(f"\n---\n# {file_name}\n{text.strip()}")
                st.success(f"✅ Parsed: {file_name}")
            else:
                st.warning(f"⚠️ No usable text in: {file_name}")

        except Exception as e:
            st.error(f"❌ Error processing {file_name}: {e}")

# === Display and Store ===
combined = "\n\n".join(all_texts)
st.session_state["persona_input_text"] = combined
st.write(f"📊 Parsed content from {len(all_texts)} file(s)")
st.text_area("📄 Combined Extracted Text", value=combined[:3000], height=300)

# === GPT Persona Generator ===
GPT_SYSTEM_PROMPT = """
You are a senior marketing strategist trained in audience research and persona development.

Based on the following content from research documents, slides, and planning inputs, generate a marketing persona including:
- First name (fictional)
- Age range
- Occupation/title
- Demographics (location, gender if inferred)
- Psychographics (attitudes, beliefs, motivations)
- Media habits
- Buying or decision behavior

Return the result as a structured persona profile using bullet points.
"""

if st.button("🤖 Generate Persona with GPT"):
    if not st.session_state.get("persona_input_text"):
        st.warning("No input text available.")
    else:
        with st.spinner("Generating persona..."):
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": GPT_SYSTEM_PROMPT},
                        {"role": "user", "content": st.session_state["persona_input_text"][:8000]}
                    ],
                    temperature=0.7,
                    max_tokens=800
                )
                output = response["choices"][0]["message"]["content"]
                st.markdown("### 🧠 Generated Persona")
                st.markdown(output)
            except Exception as e:
                st.error(f"GPT error: {e}")



